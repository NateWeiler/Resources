from pptx import Presentation
import os


def create_presentation(filename, content):
    "Create a presentation using a multiline string"
    # transform the multiline in a list
    content2 = []
    for line in content.split("\n\n"):
        line = line.split("\n")
        content2.append(line)
    print(content2)
    prs = Presentation()
    layout = prs.slide_layouts[1]
    for lst in content2:
        slide = prs.slides
        slide = slide.add_slide(layout)
        slide.shapes.title.text = lst[0]
        body = slide.placeholders[1]
        tf = body.text_frame
        if "-" in lst[1]:
            for item in lst[1].split("-")[1:]:
                p = tf.add_paragraph()
                p.text = item
                p.level = 1
        else:
            tf.text = lst[1]
    # save and launch the file
    prs.save(filename)
    os.startfile(filename)


content = """Chapter 1
This is the start

Chapter 2
Important things: - happiness - money - friends - coding"""

create_presentation("example.pptx", content)