import tkinter as tk


content = """Esercizio sul business plan
Rispondi alle seguenti domande
1. Cos'è un business plan?
2. Da quali parti è composto un business plan?
3. A chi si rivolge il business plan?
4. Cosa analizza la Swot Analysis?
5. Che rapporto c'è tra paino di marketing e business plan?
Analisi di un caso concreto
Due giovani neolaureati vogliono investire nell’area in cui sono cresciuti: vivono in un territorio ad alta attrattiva turistica, che confina con una zona balneare e con una zona montana.
Tre amici, laureati in Economia aziendale, vogliono realizzare un'impresa nel loro territorio situato in una posizione che potrebbe attrarre molti turisti per la sua posizione: si trova al centro di una zona che comprende, ai suoi estremi, il mare e la montagna a pochi km di distanza.
Uno dei tre amici possiede un immobile che potrebbe essere ristrutturato per potervi avviare l'attività: un ristorante con annesso Hotel. La loro intenzione p quella di rivolgersi ad una clientela giovane, come loro, che ama stare a contatto con la natura e che ama la cucina sana locale.
Da un'indagine di mercato sulla quale i futuri soci hanno basato le loro ricerche, risulta che la maggior parte dei potenziali clienti ritiene che è necessario migliorare l'accoglienza turistica, in quanto non sono presenti strutture che soddisfino le loro aspettative. Ci si affida all'affitto di case da private che non hanno una capacità di gestire la clientela in maniera professionale e non hanno il tempo e la capacità di offrire quella rapidità e attenzione al cliente che oggi è richiesta dalla clientela di fascia di età tra i 20 e i 35 anni, oggetto di interesse per l'impresa da costituire.
Dall'analisi risulta che oltre il 40% dei clienti vorrebbe potersi informare sulle offerte turistiche tramite il web. Una buona parte, inoltre, amerebbe poter verificare tramite il web di occasioni e offerte promozionali. Molti di loro hanno l'abitudine di consultare prima le recensioni per verificare l'affidabilità delle strutture ricettive e la soddisfazione dei clienti.
Il 50% dei clienti ritiene, inoltre, che occorre migliorare le proposte riguardanti itinerari del territorio, in quanto sono  interessati anche a conoscere i prodotti enogastronomici e il folklore locale. I turisti si fermano in media per 3 giorni. Il periodo di maggiore afflusso turistico si estende da aprile a settembre.
End..."""

counter = 0


def next():
    global counter, list_content
    if counter < len(list_content) -1:
        counter += 1
    text.delete("1.0", tk.END)
    text.insert("1.0", list_content[counter])
    button["text"] = "Next " + str(counter) + "/" + str(len(list_content)-1)


def back():
    global counter, list_content
    if counter > 0:
        counter -= 1
    text.delete("1.0", tk.END)
    text.insert("1.0", list_content[counter])
    button["text"] = "Next " + str(counter) + "/" + str(len(list_content)-1)

root = tk.Tk()
text = tk.Text(root, wrap=tk.WORD, font="Arial 20", width=50, height=15)
text.pack(fill='both')
list_content = content.split("\n")
button = tk.Button(
    root, text="Next " + str(counter), command=next, font="Arial 20")
button.pack()
button2 = tk.Button(
    root, text="back ", command=back, font="Arial 20")
button2.pack()
root.mainloop()

/*
Extract text from powerpoint (to switch to presenter)
from pptx import Presentation
import glob
*/

for eachfile in glob.glob("*.pptx"):
    prs = Presentation(eachfile)
    print(eachfile)
    print("----------------------")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)